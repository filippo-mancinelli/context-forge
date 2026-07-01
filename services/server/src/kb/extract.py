"""Text extraction from uploaded knowledge-base documents.

Supports a wide range of formats. Each extractor imports its (optional)
dependency lazily so the base install stays lean: if a library or the system
``tesseract`` binary is missing, the caller receives a clear, actionable
``ExtractionError`` that surfaces to the document's ``error_message``.

Supported formats
-----------------
- PDF                        → pypdf
- Word (.docx)               → python-docx
- Excel (.xlsx/.xlsm)        → openpyxl
- Excel legacy (.xls)        → xlrd
- PowerPoint (.pptx)         → python-pptx
- Rich text (.rtf)           → striprtf
- HTML (.html/.htm)          → BeautifulSoup (falls back to a naive tag strip)
- Images (png/jpg/…)         → OCR via pytesseract + Pillow
- CSV / TSV                  → decoded text
- Plain text, Markdown, code, JSON, YAML, XML, …  → decoded text
"""
from __future__ import annotations

import csv
import io
import json
import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)


class ExtractionError(Exception):
    """Raised when a document cannot be converted to text."""


@dataclass
class ExtractedDocument:
    text: str
    metadata: dict[str, Any] = field(default_factory=dict)


# Extensions we treat as UTF-8-ish text and read directly.
TEXT_EXTENSIONS = {
    ".txt", ".text", ".md", ".markdown", ".rst", ".log", ".rtfd",
    ".json", ".jsonl", ".ndjson", ".yaml", ".yml", ".toml", ".ini",
    ".cfg", ".conf", ".env", ".properties", ".xml", ".csv", ".tsv",
    # Source code — indexing arbitrary snippets is useful in a knowledge base.
    ".py", ".js", ".jsx", ".ts", ".tsx", ".go", ".java", ".rb", ".rs",
    ".c", ".h", ".cpp", ".hpp", ".cc", ".cs", ".php", ".swift", ".kt",
    ".scala", ".sh", ".bash", ".zsh", ".sql", ".css", ".scss", ".less",
    ".html", ".htm", ".vue", ".svelte", ".dockerfile", ".proto", ".graphql",
}

IMAGE_EXTENSIONS = {
    ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tif", ".tiff", ".webp",
}

# Every extension the uploader accepts (used for client-side hints & validation).
SUPPORTED_EXTENSIONS = (
    TEXT_EXTENSIONS
    | IMAGE_EXTENSIONS
    | {".pdf", ".docx", ".doc", ".xlsx", ".xlsm", ".xls", ".pptx", ".rtf"}
)


def is_supported(filename: str) -> bool:
    return Path(filename).suffix.lower() in SUPPORTED_EXTENSIONS


def extract_text(path: str | Path, *, filename: str | None = None) -> ExtractedDocument:
    """Extract plain text from a file on disk.

    ``filename`` overrides the extension detection when the stored path has a
    generated name (e.g. ``12__report.pdf``).
    """
    p = Path(path)
    name = filename or p.name
    ext = Path(name).suffix.lower()

    try:
        if ext == ".pdf":
            return _extract_pdf(p)
        if ext == ".docx":
            return _extract_docx(p)
        if ext == ".doc":
            return _extract_legacy_doc(p)
        if ext in (".xlsx", ".xlsm"):
            return _extract_xlsx(p)
        if ext == ".xls":
            return _extract_xls(p)
        if ext == ".pptx":
            return _extract_pptx(p)
        if ext == ".rtf":
            return _extract_rtf(p)
        if ext in (".html", ".htm"):
            return _extract_html(p)
        if ext in IMAGE_EXTENSIONS:
            return _extract_image_ocr(p)
        if ext in (".csv", ".tsv"):
            return _extract_delimited(p, ext)
        if ext in TEXT_EXTENSIONS or ext == "":
            return _extract_plain_text(p)
    except ExtractionError:
        raise
    except Exception as e:  # noqa: BLE001 — normalize any lib error
        raise ExtractionError(f"Failed to extract '{name}': {e}") from e

    raise ExtractionError(f"Unsupported file type: '{ext or name}'")


# --------------------------------------------------------------------------- #
# Format-specific extractors
# --------------------------------------------------------------------------- #

def _read_bytes(p: Path) -> bytes:
    return p.read_bytes()


def _decode_bytes(data: bytes) -> str:
    for encoding in ("utf-8", "utf-16", "latin-1"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _extract_plain_text(p: Path) -> ExtractedDocument:
    return ExtractedDocument(text=_decode_bytes(_read_bytes(p)))


def _extract_delimited(p: Path, ext: str) -> ExtractedDocument:
    raw = _decode_bytes(_read_bytes(p))
    delimiter = "\t" if ext == ".tsv" else ","
    lines: list[str] = []
    try:
        reader = csv.reader(io.StringIO(raw), delimiter=delimiter)
        for row in reader:
            if any(cell.strip() for cell in row):
                lines.append(" | ".join(cell.strip() for cell in row))
    except Exception:
        return ExtractedDocument(text=raw)
    return ExtractedDocument(text="\n".join(lines), metadata={"rows": len(lines)})


def _extract_pdf(p: Path) -> ExtractedDocument:
    try:
        from pypdf import PdfReader
    except ImportError as e:
        raise ExtractionError(
            "PDF support requires the 'pypdf' package to be installed."
        ) from e

    reader = PdfReader(str(p))
    if getattr(reader, "is_encrypted", False):
        try:
            reader.decrypt("")
        except Exception:
            raise ExtractionError("PDF is encrypted and cannot be read.")

    parts: list[str] = []
    for page in reader.pages:
        try:
            parts.append(page.extract_text() or "")
        except Exception:
            parts.append("")
    text = "\n\n".join(parts).strip()
    meta = {"pages": len(reader.pages)}
    if not text:
        raise ExtractionError(
            "No extractable text in this PDF (it may be a scanned document — "
            "upload it as an image for OCR)."
        )
    return ExtractedDocument(text=text, metadata=meta)


def _extract_docx(p: Path) -> ExtractedDocument:
    try:
        import docx  # python-docx
    except ImportError as e:
        raise ExtractionError(
            "Word (.docx) support requires the 'python-docx' package."
        ) from e

    document = docx.Document(str(p))
    parts = [para.text for para in document.paragraphs if para.text.strip()]

    # Include table cell text — often where the real content lives.
    for table in document.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    text = "\n".join(parts).strip()
    if not text:
        raise ExtractionError("The document contains no extractable text.")
    return ExtractedDocument(text=text, metadata={"paragraphs": len(parts)})


def _extract_legacy_doc(p: Path) -> ExtractedDocument:
    # Legacy binary .doc needs antiword/LibreOffice which we don't bundle.
    # Try python-docx in case it's actually a mislabeled .docx, else guide the user.
    try:
        return _extract_docx(p)
    except Exception:
        raise ExtractionError(
            "Legacy '.doc' files are not supported directly. "
            "Please convert the file to .docx (or PDF) and re-upload."
        )


def _extract_xlsx(p: Path) -> ExtractedDocument:
    try:
        from openpyxl import load_workbook
    except ImportError as e:
        raise ExtractionError(
            "Excel (.xlsx) support requires the 'openpyxl' package."
        ) from e

    wb = load_workbook(filename=str(p), read_only=True, data_only=True)
    parts: list[str] = []
    for ws in wb.worksheets:
        parts.append(f"# Sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None and str(c).strip()]
            if cells:
                parts.append(" | ".join(cells))
    wb.close()
    text = "\n".join(parts).strip()
    if not text:
        raise ExtractionError("The spreadsheet contains no data.")
    return ExtractedDocument(text=text, metadata={"sheets": len(wb.sheetnames)})


def _extract_xls(p: Path) -> ExtractedDocument:
    try:
        import xlrd
    except ImportError as e:
        raise ExtractionError(
            "Legacy Excel (.xls) support requires the 'xlrd' package."
        ) from e

    book = xlrd.open_workbook(str(p))
    parts: list[str] = []
    for sheet in book.sheets():
        parts.append(f"# Sheet: {sheet.name}")
        for r in range(sheet.nrows):
            cells = [str(c) for c in sheet.row_values(r) if str(c).strip()]
            if cells:
                parts.append(" | ".join(cells))
    text = "\n".join(parts).strip()
    if not text:
        raise ExtractionError("The spreadsheet contains no data.")
    return ExtractedDocument(text=text, metadata={"sheets": book.nsheets})


def _extract_pptx(p: Path) -> ExtractedDocument:
    try:
        from pptx import Presentation
    except ImportError as e:
        raise ExtractionError(
            "PowerPoint (.pptx) support requires the 'python-pptx' package."
        ) from e

    prs = Presentation(str(p))
    parts: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        slide_parts.append(line)
        if slide_parts:
            parts.append(f"# Slide {idx}")
            parts.extend(slide_parts)
    text = "\n".join(parts).strip()
    if not text:
        raise ExtractionError("The presentation contains no extractable text.")
    return ExtractedDocument(text=text, metadata={"slides": len(prs.slides._sldIdLst)})


def _extract_rtf(p: Path) -> ExtractedDocument:
    raw = _decode_bytes(_read_bytes(p))
    try:
        from striprtf.striprtf import rtf_to_text

        text = rtf_to_text(raw)
    except ImportError:
        raise ExtractionError(
            "Rich text (.rtf) support requires the 'striprtf' package."
        )
    text = (text or "").strip()
    if not text:
        raise ExtractionError("The document contains no extractable text.")
    return ExtractedDocument(text=text)


def _extract_html(p: Path) -> ExtractedDocument:
    raw = _decode_bytes(_read_bytes(p))
    try:
        from bs4 import BeautifulSoup

        soup = BeautifulSoup(raw, "html.parser")
        for tag in soup(["script", "style", "noscript"]):
            tag.decompose()
        text = soup.get_text(separator="\n")
    except ImportError:
        # Fallback: naive tag strip so HTML still works without bs4.
        import re

        text = re.sub(r"<[^>]+>", " ", raw)
    lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
    joined = "\n".join(lines).strip()
    if not joined:
        raise ExtractionError("The HTML document contains no readable text.")
    return ExtractedDocument(text=joined)


def _extract_image_ocr(p: Path) -> ExtractedDocument:
    try:
        import pytesseract
        from PIL import Image
    except ImportError as e:
        raise ExtractionError(
            "Image OCR requires the 'pytesseract' and 'Pillow' packages."
        ) from e

    try:
        image = Image.open(str(p))
    except Exception as e:
        raise ExtractionError(f"Could not open image: {e}") from e

    try:
        text = pytesseract.image_to_string(image)
    except pytesseract.TesseractNotFoundError as e:
        raise ExtractionError(
            "OCR requires the 'tesseract' binary to be installed on the server."
        ) from e
    except Exception as e:
        raise ExtractionError(f"OCR failed: {e}") from e

    text = (text or "").strip()
    meta = {"width": image.width, "height": image.height, "ocr": True}
    if not text:
        raise ExtractionError(
            "No text could be recognized in this image via OCR."
        )
    return ExtractedDocument(text=text, metadata=meta)
